# -*- coding: utf-8 -*-
"""
Cafe Suzak — Content Marketing Progress Report
B&W base + Crimson accent + Cream beige, Gmarket Sans
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ---------- Palette ----------
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x26, 0x26, 0x26)   # 262626
INK_2   = RGBColor(0x55, 0x55, 0x55)   # secondary text
INK_3   = RGBColor(0x99, 0x99, 0x99)   # caption / muted
CREAM   = RGBColor(0xFA, 0xE7, 0xCB)   # FAE7CB
CREAM_2 = RGBColor(0xFD, 0xF3, 0xE2)   # softer wash
CRIMSON = RGBColor(0xDC, 0x14, 0x3C)   # crimson point
LINE    = RGBColor(0xE5, 0xE5, 0xE5)

# ---------- Fonts ----------
F_SANS  = "Gmarket Sans"
F_SANS_LT = "Gmarket Sans Light"
F_SANS_BD = "Gmarket Sans Bold"

# ---------- Slide size: 16:9 widescreen ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ---------- Helpers ----------
def add_slide():
    s = prs.slides.add_slide(BLANK)
    # white background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.shadow.inherit = False
    return s

def add_rect(s, x, y, w, h, fill=None, line=None, line_w=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        if line_w is not None:
            sh.line.width = line_w
    sh.shadow.inherit = False
    return sh

def add_line(s, x1, y1, x2, y2, color=INK, weight=0.75):
    ln = s.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln

def add_text(s, x, y, w, h, text, *,
             font=F_SANS, size=12, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
             tracking=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if tracking is not None:
            # letter spacing via XML
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(tracking))  # in 100ths of point
    return tb

def add_text_runs(s, x, y, w, h, runs, *, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, line_spacing=1.15):
    """runs = list of (text, dict(opts))"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, opts in runs:
        if text == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            continue
        r = p.add_run()
        r.text = text
        r.font.name = opts.get("font", F_SANS)
        r.font.size = Pt(opts.get("size", 12))
        r.font.bold = opts.get("bold", False)
        r.font.color.rgb = opts.get("color", INK)
        if opts.get("tracking") is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(opts["tracking"]))
    return tb

# ---------- Page header (chapter mark) ----------
def page_header(s, idx, total, chapter_kr, chapter_en):
    # left top: index + chapter
    add_text(s, Inches(0.6), Inches(0.45), Inches(7), Inches(0.3),
             f"{idx:02d} / {total:02d}", size=9, color=CRIMSON, bold=True, tracking=300)
    add_text(s, Inches(1.2), Inches(0.45), Inches(7), Inches(0.3),
             f"  {chapter_en}", size=9, color=INK_3, tracking=200)
    # right top: brand
    add_text(s, Inches(11), Inches(0.45), Inches(2), Inches(0.3),
             "CAFE SUZAK · 2026", size=9, color=INK_3, align=PP_ALIGN.RIGHT, tracking=200)
    # thin underline across top
    add_line(s, Inches(0.6), Inches(0.78), Inches(12.733), Inches(0.78), color=LINE, weight=0.5)

def page_footer(s, page_label):
    add_line(s, Inches(0.6), Inches(7.1), Inches(12.733), Inches(7.1), color=LINE, weight=0.5)
    add_text(s, Inches(0.6), Inches(7.15), Inches(6), Inches(0.25),
             "글로벌경영학과  2022390786  안혜빈", size=8, color=INK_3, tracking=200)
    add_text(s, Inches(7), Inches(7.15), Inches(5.7), Inches(0.25),
             page_label, size=8, color=INK_3, align=PP_ALIGN.RIGHT, tracking=200)

# ===================================================================
# SLIDE 1 — Cover
# ===================================================================
s = add_slide()

# left vertical cream sidebar
add_rect(s, 0, 0, Inches(4.4), SH, fill=CREAM)

# crimson short line on cover
add_rect(s, Inches(0.8), Inches(1.0), Inches(0.6), Inches(0.05), fill=CRIMSON)

# Sidebar text - small label
add_text(s, Inches(0.8), Inches(1.25), Inches(3.5), Inches(0.4),
         "CONTENT MARKETING\nPROGRESS REPORT", size=10, color=INK, bold=True, tracking=400, line_spacing=1.4)

# Sidebar main title
add_text_runs(s, Inches(0.8), Inches(5.4), Inches(3.5), Inches(1.4),
              [("CAFE\nSUZAK.", {"size":36,"bold":True,"color":INK,"tracking":-20})],
              line_spacing=1.0)

# Sidebar footer
add_text(s, Inches(0.8), Inches(6.8), Inches(3.5), Inches(0.3),
         "Sejong · Jochiwon · Roastery", size=9, color=INK_2, tracking=300)

# Right main area
RX = Inches(5.0)
add_text(s, RX, Inches(1.1), Inches(7.5), Inches(0.4),
         "세종충정지역 영세기업 AI 콘텐츠 마케팅 산학협력",
         size=11, color=INK_2, tracking=200)

# Big headline
add_text_runs(s, RX, Inches(1.6), Inches(7.7), Inches(2.6),
              [("조용히 쌓아온 가치를,\n", {"size":36,"bold":True,"color":INK,"tracking":-20}),
               ("새로운 채널로 ", {"size":36,"bold":True,"color":INK,"tracking":-20}),
               ("닿다", {"size":36,"bold":True,"color":CRIMSON,"tracking":-20}),
               (".", {"size":36,"bold":True,"color":INK,"tracking":-20})],
              line_spacing=1.15)

# Subhead
add_text(s, RX, Inches(4.3), Inches(7.5), Inches(0.5),
         "카페수작 · Content Marketing Mid-term Report",
         size=14, color=INK, bold=True, tracking=100)

# divider
add_rect(s, RX, Inches(4.85), Inches(7.5), Inches(0.01), fill=INK)

# meta
add_text_runs(s, RX, Inches(5.05), Inches(7.5), Inches(0.4),
              [("CLIENT", {"size":9,"color":INK_3,"tracking":300,"bold":True})])
add_text(s, RX, Inches(5.3), Inches(7.5), Inches(0.4),
         "카페수작 (CAFE SUZAK) · 세종 조치원읍 원마루길 28-1",
         size=11, color=INK)

add_text_runs(s, RX, Inches(5.85), Inches(7.5), Inches(0.4),
              [("PROJECT", {"size":9,"color":INK_3,"tracking":300,"bold":True})])
add_text(s, RX, Inches(6.1), Inches(7.5), Inches(0.4),
         "AI 콘텐츠 기반 다채널 마케팅 캠페인",
         size=11, color=INK)

add_text_runs(s, RX, Inches(6.65), Inches(7.5), Inches(0.4),
              [("AUTHOR  ·  DATE", {"size":9,"color":INK_3,"tracking":300,"bold":True})])
add_text(s, RX, Inches(6.9), Inches(7.5), Inches(0.4),
         "글로벌경영학과 2022390786 안혜빈  ·  2026. 05",
         size=11, color=INK)

# ===================================================================
# SLIDE 2 — Index
# ===================================================================
s = add_slide()
page_header(s, 0, 12, "목차", "INDEX")

add_text(s, Inches(0.6), Inches(1.3), Inches(8), Inches(0.5),
         "INDEX", size=11, color=CRIMSON, bold=True, tracking=400)
add_text_runs(s, Inches(0.6), Inches(1.7), Inches(10), Inches(1.5),
              [("이 보고서는 ", {"size":30,"color":INK,"bold":True,"tracking":-20}),
               ("어떻게 카페수작의 가치를\n", {"size":30,"color":INK,"bold":True,"tracking":-20}),
               ("새로운 채널에 ", {"size":30,"color":INK,"bold":True,"tracking":-20}),
               ("닿게 했는지", {"size":30,"color":CRIMSON,"bold":True,"tracking":-20}),
               (" 정리합니다.", {"size":30,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.25)

# Index list — 2 columns
items_left = [
    ("01", "프로젝트 배경", "Project Background"),
    ("02", "문제 정의 — 가치 vs 인지도", "The Paradox"),
    ("03", "타겟 고객", "Target Audience"),
    ("04", "마케팅 전략 개요", "Strategy Overview"),
    ("05", "AI 활용 제작 프로세스", "AI-Driven Process"),
    ("06", "산출물 · 웹페이지", "Output · Webpage"),
]
items_right = [
    ("07", "산출물 · 블로그", "Output · Blog"),
    ("08", "산출물 · 인스타그램", "Output · Instagram"),
    ("09", "산출물 · 에브리타임", "Output · Everytime"),
    ("10", "유통 채널 전략", "Channel Strategy"),
    ("11", "성과 측정 프레임", "Performance Framework"),
    ("12", "향후 계획", "Next Step"),
]

def index_row(s, x, y, num, kr, en):
    add_text(s, x, y, Inches(0.6), Inches(0.4),
             num, size=11, color=CRIMSON, bold=True, tracking=200)
    add_text(s, x+Inches(0.7), y-Inches(0.02), Inches(5), Inches(0.3),
             kr, size=15, color=INK, bold=True)
    add_text(s, x+Inches(0.7), y+Inches(0.3), Inches(5), Inches(0.25),
             en, size=9, color=INK_3, tracking=200)
    add_line(s, x, y+Inches(0.65), x+Inches(5.5), y+Inches(0.65), color=LINE, weight=0.5)

y0 = Inches(4.0)
for i, (n, kr, en) in enumerate(items_left):
    index_row(s, Inches(0.6), y0+Inches(0.5*i), n, kr, en)
for i, (n, kr, en) in enumerate(items_right):
    index_row(s, Inches(7.0), y0+Inches(0.5*i), n, kr, en)

page_footer(s, "INDEX")

# ===================================================================
# SLIDE 3 — Project Background
# ===================================================================
s = add_slide()
page_header(s, 1, 12, "프로젝트 배경", "PROJECT BACKGROUND")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "01", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(11), Inches(1.5),
              [("진짜 커피를 아는 사람들이 ", {"size":32,"color":INK,"bold":True,"tracking":-20}),
               ("결국 찾는 곳", {"size":32,"color":CRIMSON,"bold":True,"tracking":-20}),
               (",\n카페 수작.", {"size":32,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.2)

add_text(s, Inches(0.6), Inches(3.4), Inches(7.5), Inches(1.2),
         "세종 조치원의 작은 골목, 원마루길 28-1.\n2012년부터 — 15년차 마스터가 매일 매장에서 직접 원두를 볶고,\n시럽·소스·디저트까지 100% 손으로 만드는 로스터리 카페입니다.",
         size=13, color=INK_2, line_spacing=1.7)

# Three core values - boxes
def value_box(s, x, y, w, h, label, kr, en):
    add_rect(s, x, y, w, h, fill=CREAM_2)
    add_rect(s, x, y, Inches(0.05), h, fill=CRIMSON)  # crimson tab
    add_text(s, x+Inches(0.3), y+Inches(0.25), w-Inches(0.4), Inches(0.3),
             label, size=9, color=CRIMSON, bold=True, tracking=300)
    add_text(s, x+Inches(0.3), y+Inches(0.55), w-Inches(0.4), Inches(0.5),
             kr, size=16, color=INK, bold=True)
    add_text(s, x+Inches(0.3), y+Inches(1.1), w-Inches(0.4), Inches(0.8),
             en, size=10, color=INK_2, line_spacing=1.55)

bx, by = Inches(0.6), Inches(5.0)
bw, bh = Inches(4.0), Inches(1.85)
value_box(s, bx,            by, bw, bh, "HANDMADE",
          "100% 수작업의 철학",
          "생두 직접 구매 · 매장 로스팅 · 시럽/소스 수제 · 파스타/리조또까지")
value_box(s, bx+Inches(4.1), by, bw, bh, "MASTERY",
          "15년차 마스터 (2012~)",
          "본질에 집착하는 결과물 · 단순 카페가 아닌 완성도 높은 커피 경험")
value_box(s, bx+Inches(8.2), by, bw, bh, "EXPERIENCE",
          "낮엔 카페, 밤엔 바",
          "아메리카노부터 모히또까지 · 22시까지 운영 · 조치원의 작은 아지트")

page_footer(s, "PROJECT BACKGROUND")

# ===================================================================
# SLIDE 4 — The Paradox (Problem)
# ===================================================================
s = add_slide()
page_header(s, 2, 12, "문제 정의", "THE PARADOX")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "02", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(11), Inches(1.5),
              [("가치는 충분하지만, ", {"size":32,"color":INK,"bold":True,"tracking":-20}),
               ("전달되지 않습니다", {"size":32,"color":CRIMSON,"bold":True,"tracking":-20}),
               (".", {"size":32,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.2)

add_text(s, Inches(0.6), Inches(3.3), Inches(11), Inches(0.5),
         "이미 완성된 명작이지만, 디지털 영역에서는 그 가치가 표현되지 못해 신규 고객 유입이 정체되어 있습니다.",
         size=12, color=INK_2, line_spacing=1.6)

# Two-column problem layout
LX, RX = Inches(0.6), Inches(6.8)
COLW = Inches(5.9)

# Left — visibility
add_rect(s, LX, Inches(4.2), COLW, Inches(2.6), fill=CREAM_2)
add_text(s, LX+Inches(0.4), Inches(4.4), COLW-Inches(0.6), Inches(0.3),
         "PERCEIVED VISIBILITY", size=9, color=CRIMSON, bold=True, tracking=300)
add_text(s, LX+Inches(0.4), Inches(4.7), COLW-Inches(0.6), Inches(0.4),
         "가벼운 인지도", size=18, color=INK, bold=True)
add_text(s, LX+Inches(0.4), Inches(5.25), COLW-Inches(0.6), Inches(1.6),
         "·  간판이 잘 보이지 않아 지나치기 쉬운 위치\n·  테이크아웃·운영 정보 등 기본 정보 부재\n·  체계적인 홍보 경험이 거의 없음",
         size=12, color=INK_2, line_spacing=1.85)

# Right — value
add_rect(s, RX, Inches(4.2), COLW, Inches(2.6), fill=INK)
add_text(s, RX+Inches(0.4), Inches(4.4), COLW-Inches(0.6), Inches(0.3),
         "INTRINSIC VALUE", size=9, color=CREAM, bold=True, tracking=300)
add_text(s, RX+Inches(0.4), Inches(4.7), COLW-Inches(0.6), Inches(0.4),
         "본질적 경쟁력", size=18, color=WHITE, bold=True)
add_text(s, RX+Inches(0.4), Inches(5.25), COLW-Inches(0.6), Inches(1.6),
         "·  2012년부터 — 15년차 탄탄한 업력\n·  세종시 최고 수준이라 자부하는 직접 로스팅 커피\n·  제빵 자격증 · 파스타/소스까지 직접 만드는 요리 역량",
         size=12, color=WHITE, line_spacing=1.85, font=F_SANS)

page_footer(s, "THE PARADOX")

# ===================================================================
# SLIDE 5 — Target Audience
# ===================================================================
s = add_slide()
page_header(s, 3, 12, "타겟 고객", "TARGET AUDIENCE")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "03", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(11), Inches(1.5),
              [("누구에게 ", {"size":32,"color":INK,"bold":True,"tracking":-20}),
               ("우리의 가치", {"size":32,"color":CRIMSON,"bold":True,"tracking":-20}),
               ("를 전달할 것인가?", {"size":32,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.2)

# 3 personas
def persona(s, x, y, w, h, tag, en, kr, desc):
    add_rect(s, x, y, w, h, fill=WHITE, line=LINE, line_w=Pt(0.75))
    add_rect(s, x, y, w, Inches(0.5), fill=INK)
    add_text(s, x+Inches(0.3), y+Inches(0.13), w-Inches(0.6), Inches(0.3),
             tag, size=10, color=CREAM, bold=True, tracking=400)
    add_text(s, x+Inches(0.3), y+Inches(0.75), w-Inches(0.6), Inches(0.4),
             en, size=9, color=CRIMSON, bold=True, tracking=300)
    add_text(s, x+Inches(0.3), y+Inches(1.05), w-Inches(0.6), Inches(0.5),
             kr, size=18, color=INK, bold=True, line_spacing=1.2)
    add_text(s, x+Inches(0.3), y+Inches(2.0), w-Inches(0.6), Inches(1.5),
             desc, size=11, color=INK_2, line_spacing=1.7)
    add_rect(s, x+Inches(0.3), y+Inches(1.85), Inches(0.5), Inches(0.02), fill=CRIMSON)

py, ph = Inches(3.3), Inches(3.5)
persona(s, Inches(0.6),  py, Inches(4.0), ph, "01  PRIMARY",
        "PRIMARY",
        "고려대 및 인근\n대학생",
        "주 고객층이자 잠재적 바이럴의 핵심 진원지.\n조치원 내 학생 커뮤니티가 좁고\n구전 효과가 크게 작용하는 지역.")
persona(s, Inches(4.8),  py, Inches(4.0), ph, "02  BEHAVIORAL",
        "BEHAVIORAL",
        "카공족 &\n장기 체류자",
        "조용한 환경과 사장님의 배려(터치 없음),\n넓은 좌석 간격으로 인해 \n장시간 공부하기 편한 공간을 찾는 고객.")
persona(s, Inches(9.0),  py, Inches(4.0), ph, "03  LOYALTY",
        "LOYALTY",
        "졸업생 &\n지인 추천 방문객",
        "오직 '커피 맛' 하나로 재방문하는\n두터운 충성도를 가진 그룹.\n선배–후배 추천 사이클의 시작점.")

page_footer(s, "TARGET AUDIENCE")

# ===================================================================
# SLIDE 6 — Strategy Overview
# ===================================================================
s = add_slide()
page_header(s, 4, 12, "마케팅 전략 개요", "STRATEGY OVERVIEW")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "04", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.5),
              [("맛으로 승부하는 시대는 끝났습니다.\n", {"size":28,"color":INK,"bold":True,"tracking":-20}),
               ("이제 ", {"size":28,"color":INK,"bold":True,"tracking":-20}),
               ("경험을 증명하고 전파", {"size":28,"color":CRIMSON,"bold":True,"tracking":-20}),
               ("해야 합니다.", {"size":28,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.25)

# Strategy pillars
add_text(s, Inches(0.6), Inches(4.0), Inches(11), Inches(0.4),
         "STRATEGIC PILLARS", size=10, color=CRIMSON, bold=True, tracking=400)
add_rect(s, Inches(0.6), Inches(4.35), Inches(12.13), Inches(0.01), fill=INK)

def pillar(s, x, y, w, num, title, body):
    add_text(s, x, y, w, Inches(0.35),
             num, size=11, color=CRIMSON, bold=True, tracking=200)
    add_text(s, x, y+Inches(0.4), w, Inches(0.4),
             title, size=15, color=INK, bold=True)
    add_text(s, x, y+Inches(0.95), w, Inches(1.5),
             body, size=11, color=INK_2, line_spacing=1.65)

py = Inches(4.6)
pillar(s, Inches(0.6),  py, Inches(3.9), "P · 01", "장인정신 스토리텔링",
       "로스팅 과정·메뉴 비하인드·공간 분위기를\n에디토리얼 톤으로 풀어 전문성과 정성을 전달.")
pillar(s, Inches(4.7),  py, Inches(3.9), "P · 02", "목적형 검색 타겟팅",
       "'조치원 카페수작', '핸드드립', '카공 카페' 등\n검색 의도가 명확한 키워드로 신뢰형 유입 확보.")
pillar(s, Inches(8.8),  py, Inches(3.9), "P · 03", "AI 콘텐츠 자산화",
       "Claude + VSCode로 웹페이지·블로그·SNS를\n빠르게 제작·반복 개선, 가성비 높은 자산으로 축적.")

page_footer(s, "STRATEGY OVERVIEW")

# ===================================================================
# SLIDE 7 — AI-Driven Process
# ===================================================================
s = add_slide()
page_header(s, 5, 12, "AI 활용 제작 프로세스", "AI-DRIVEN PROCESS")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "05", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.5),
              [("Claude Code + VS Code", {"size":32,"color":INK,"bold":True,"tracking":-20}),
               ("로\n", {"size":32,"color":INK,"bold":True,"tracking":-20}),
               ("아이디어를 실제 자산으로", {"size":32,"color":CRIMSON,"bold":True,"tracking":-20}),
               (".", {"size":32,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.2)

# Tool blocks
add_text(s, Inches(0.6), Inches(3.7), Inches(11), Inches(0.4),
         "TOOLSTACK", size=10, color=CRIMSON, bold=True, tracking=400)
add_rect(s, Inches(0.6), Inches(4.05), Inches(12.13), Inches(0.01), fill=INK)

def tool(s, x, y, w, name, role, detail):
    add_rect(s, x, y, w, Inches(2.2), fill=CREAM_2)
    add_text(s, x+Inches(0.3), y+Inches(0.25), w-Inches(0.6), Inches(0.3),
             name, size=14, color=INK, bold=True)
    add_text(s, x+Inches(0.3), y+Inches(0.7), w-Inches(0.6), Inches(0.3),
             role, size=10, color=CRIMSON, bold=True, tracking=200)
    add_text(s, x+Inches(0.3), y+Inches(1.05), w-Inches(0.6), Inches(1.1),
             detail, size=10, color=INK_2, line_spacing=1.7)

py = Inches(4.3)
tool(s, Inches(0.6), py, Inches(3.9), "Claude Code",
     "AI CO-WORKER",
     "프롬프트로 코드·카피·디자인 톤·문서를 즉시 생성.\n반복 수정에 강해 빠른 디자인 이터레이션 가능.")
tool(s, Inches(4.7), py, Inches(3.9), "VS Code",
     "EDITOR",
     "코드 결과물 확인·미세 조정·로컬 프리뷰.\n실시간 미리보기로 톤·반응형 검수 진행.")
tool(s, Inches(8.8), py, Inches(3.9), "Netlify Drop",
     "DEPLOY",
     "별도 계정 없이 폴더 드래그만으로\n공개 URL 즉시 발급. 모바일 테스트 즉시 가능.")

# Process steps below
add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.4),
         "PROMPT → ITERATE → DEPLOY  ·  약 15+ 차례 반복 수정으로 콘셉트·톤·기능을 점진적으로 다듬어 완성",
         size=10, color=INK_2, tracking=100)

page_footer(s, "AI-DRIVEN PROCESS")

# ===================================================================
# SLIDE 8 — Output 1: Webpage
# ===================================================================
s = add_slide()
page_header(s, 6, 12, "산출물 · 웹페이지", "OUTPUT · WEBPAGE")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "06", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(11), Inches(1.6),
              [("\"오늘 카페수작에서 ", {"size":30,"color":INK,"bold":True,"tracking":-20}),
               ("뭐 마실까", {"size":30,"color":CRIMSON,"bold":True,"tracking":-20}),
               ("?\"", {"size":30,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.2)
add_text(s, Inches(0.6), Inches(2.65), Inches(11), Inches(0.5),
         "기분과 날씨에 맞춰 카페수작의 메뉴를 추천해주는 인터랙티브 랜딩페이지",
         size=12, color=INK_2, line_spacing=1.5)

# URL bar
add_rect(s, Inches(0.6), Inches(3.4), Inches(12.13), Inches(0.55), fill=INK)
add_text(s, Inches(0.85), Inches(3.5), Inches(2), Inches(0.35),
         "LIVE URL", size=9, color=CREAM, bold=True, tracking=300)
add_text(s, Inches(2.5), Inches(3.5), Inches(10), Inches(0.35),
         "https://meet-at-suzak.vercel.app/", size=12, color=WHITE, bold=True)

# Feature grid
add_text(s, Inches(0.6), Inches(4.2), Inches(11), Inches(0.4),
         "KEY FEATURES", size=10, color=CRIMSON, bold=True, tracking=400)
add_rect(s, Inches(0.6), Inches(4.55), Inches(12.13), Inches(0.01), fill=INK)

def feat(s, x, y, w, num, title, body):
    add_text(s, x, y, w, Inches(0.3),
             num, size=11, color=CRIMSON, bold=True, tracking=200)
    add_text(s, x, y+Inches(0.35), w, Inches(0.4),
             title, size=13, color=INK, bold=True)
    add_text(s, x, y+Inches(0.8), w, Inches(1.5),
             body, size=10, color=INK_2, line_spacing=1.7)

py = Inches(4.8)
feat(s, Inches(0.6),  py, Inches(3.0), "F · 01", "실시간 날씨 추천",
     "Open-Meteo API로 조치원 + 사용자 위치 동시 표시. 기온·강수 기반 자동 코멘트.")
feat(s, Inches(3.7),  py, Inches(3.0), "F · 02", "기분 10종 메뉴 추천",
     "잠깨야해·카공·단거·하루마무리 등 10가지 기분별 메뉴 3종 자동 매칭.")
feat(s, Inches(6.8),  py, Inches(3.0), "F · 03", "실시간 영업 상태",
     "한국시간 기준 영업/휴무 자동 판별. 토요일 휴무·라스트오더 30분 전 안내.")
feat(s, Inches(9.9),  py, Inches(2.8), "F · 04", "반응형 디자인",
     "데스크탑·태블릿·모바일 자동 대응. WatchHouse·Hola Coffee 톤 에디토리얼 무드.")

page_footer(s, "OUTPUT · WEBPAGE")

# ===================================================================
# SLIDE 9 — Output 2: Blog
# ===================================================================
s = add_slide()
page_header(s, 7, 12, "산출물 · 블로그", "OUTPUT · BLOG")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "07", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.6),
              [("전체적으로 ", {"size":28,"color":INK,"bold":True,"tracking":-20}),
               ("\"수작이라는 카페\"", {"size":28,"color":CRIMSON,"bold":True,"tracking":-20}),
               ("를\n소개하는 ", {"size":28,"color":INK,"bold":True,"tracking":-20}),
               ("마케팅 블로그", {"size":28,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.2)

# Left content
LX = Inches(0.6); LW = Inches(7.0)
add_text(s, LX, Inches(3.55), LW, Inches(0.35),
         "PLATFORM", size=9, color=CRIMSON, bold=True, tracking=300)
add_text(s, LX, Inches(3.85), LW, Inches(0.4),
         "네이버 블로그 (Naver Blog)", size=15, color=INK, bold=True)
add_text(s, LX, Inches(4.25), LW, Inches(0.5),
         "한국 검색 점유율 압도적 · 장문 + 이미지 콘텐츠에 가장 적합한 환경",
         size=11, color=INK_2, line_spacing=1.5)

add_text(s, LX, Inches(5.0), LW, Inches(0.35),
         "TONE & ANGLE", size=9, color=CRIMSON, bold=True, tracking=300)
add_text(s, LX, Inches(5.3), LW, Inches(1.5),
         "·  방문 경험 기반 1인칭 친근체로 작성\n·  카페수작의 전반적인 소개 (위치·메뉴·분위기·시그니처)\n·  글 마지막에 \"여러분도 방문해보세요\" CTA + 웹페이지 링크 자연 연결",
         size=11, color=INK_2, line_spacing=1.85)

# Right card — sample
RX = Inches(8.0); RW = Inches(4.75)
add_rect(s, RX, Inches(3.5), RW, Inches(3.3), fill=CREAM_2)
add_text(s, RX+Inches(0.35), Inches(3.7), RW-Inches(0.5), Inches(0.3),
         "SAMPLE  ·  최근 발행 글", size=9, color=CRIMSON, bold=True, tracking=300)
add_text(s, RX+Inches(0.35), Inches(4.05), RW-Inches(0.5), Inches(1.3),
         "오늘 날씨엔 뭐 마시지?\n조치원 카페수작에서\n기분별로 골라본 6잔",
         size=16, color=INK, bold=True, line_spacing=1.25)
add_rect(s, RX+Inches(0.35), Inches(5.6), Inches(0.5), Inches(0.03), fill=CRIMSON)
add_text(s, RX+Inches(0.35), Inches(5.75), RW-Inches(0.5), Inches(1.0),
         "· 약 3,000자 / 읽는 시간 5분\n· 6가지 기분 시나리오 + 메뉴 매칭\n· 마지막에 웹페이지 링크 임베드\n· 해시태그 16개 (조치원 SEO 위주)",
         size=10, color=INK_2, line_spacing=1.75)

page_footer(s, "OUTPUT · BLOG")

# ===================================================================
# SLIDE 10 — Output 3: Instagram
# ===================================================================
s = add_slide()
page_header(s, 8, 12, "산출물 · 인스타그램", "OUTPUT · INSTAGRAM")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "08", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.6),
              [("지인과 ", {"size":28,"color":INK,"bold":True,"tracking":-20}),
               ("근거리 거주자", {"size":28,"color":CRIMSON,"bold":True,"tracking":-20}),
               ("를 향한\n", {"size":28,"color":INK,"bold":True,"tracking":-20}),
               ("개인 계정 기반 자연 확산", {"size":28,"color":INK,"bold":True,"tracking":-20}),
               (".", {"size":28,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.25)

# Layout: left text, right reels card
LX = Inches(0.6); LW = Inches(7.0)
add_text(s, LX, Inches(3.7), LW, Inches(0.35),
         "PLATFORM  ·  TARGET", size=9, color=CRIMSON, bold=True, tracking=300)
add_text(s, LX, Inches(4.0), LW, Inches(0.4),
         "개인 인스타그램 (Personal IG)", size=15, color=INK, bold=True)
add_text(s, LX, Inches(4.45), LW, Inches(0.5),
         "지인 + 조치원·세종 근거리 거주자 / 학생 네트워크 1차 확산",
         size=11, color=INK_2)

add_text(s, LX, Inches(5.15), LW, Inches(0.35),
         "STRATEGY", size=9, color=CRIMSON, bold=True, tracking=300)
add_text(s, LX, Inches(5.45), LW, Inches(1.5),
         "·  본 광고가 아닌 \"내가 좋아하는 카페\" 톤의 자연스러운 추천\n·  지인 도달률이 높은 IG 특성을 활용한 신뢰 기반 확산\n·  릴스 포맷으로 짧고 강한 비주얼 메시지 전달",
         size=11, color=INK_2, line_spacing=1.85)

# Right reels card
RX = Inches(8.0); RW = Inches(4.75)
add_rect(s, RX, Inches(3.6), RW, Inches(3.3), fill=INK)
add_text(s, RX+Inches(0.35), Inches(3.8), RW-Inches(0.5), Inches(0.3),
         "DELIVERABLES", size=9, color=CREAM, bold=True, tracking=300)
add_text_runs(s, RX+Inches(0.35), Inches(4.15), RW-Inches(0.5), Inches(1.6),
              [("03", {"size":58,"color":CREAM,"bold":True,"tracking":-30}),
               ("\n", {"size":58,"color":WHITE}),
               ("Reels Uploaded", {"size":12,"color":WHITE,"bold":True,"tracking":200})],
              line_spacing=1.0)
add_rect(s, RX+Inches(0.35), Inches(5.9), Inches(0.5), Inches(0.03), fill=CRIMSON)
add_text(s, RX+Inches(0.35), Inches(6.05), RW-Inches(0.5), Inches(0.8),
         "카페 분위기 · 메뉴 클로즈업 · 핸드드립 과정\n각 릴스는 30초 내외, 음악·자막으로 몰입감 강화.",
         size=10, color=WHITE, line_spacing=1.7)

page_footer(s, "OUTPUT · INSTAGRAM")

# ===================================================================
# SLIDE 11 — Output 4: Everytime
# ===================================================================
s = add_slide()
page_header(s, 9, 12, "산출물 · 에브리타임", "OUTPUT · EVERYTIME")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "09", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.6),
              [("학생들이 ", {"size":28,"color":INK,"bold":True,"tracking":-20}),
               ("거의 매일 여는 앱", {"size":28,"color":CRIMSON,"bold":True,"tracking":-20}),
               (",\n에브리타임을 마케팅 자산화.", {"size":28,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.2)

LX = Inches(0.6); LW = Inches(7.0)
add_text(s, LX, Inches(3.65), LW, Inches(0.35),
         "PLATFORM  ·  CONTEXT", size=9, color=CRIMSON, bold=True, tracking=300)
add_text(s, LX, Inches(3.95), LW, Inches(0.4),
         "에브리타임 (Everytime · 대학 커뮤니티)", size=15, color=INK, bold=True)
add_text(s, LX, Inches(4.4), LW, Inches(0.7),
         "수업·시간표·게시판이 통합된 학생 일상 앱.\n수업이 있는 날에는 거의 매일 접속해 도달률·관여도가 압도적.",
         size=11, color=INK_2, line_spacing=1.5)

add_text(s, LX, Inches(5.3), LW, Inches(0.35),
         "STRATEGY", size=9, color=CRIMSON, bold=True, tracking=300)
add_text(s, LX, Inches(5.6), LW, Inches(1.5),
         "·  학교 인근 카페라는 강한 지리적 연결성을 부각\n·  광고 톤이 아닌 \"실제 학생 후기\" 형태로 자연스럽게 게시\n·  카공·시험 기간·데이트 등 학생 컨텍스트에 메뉴 추천 매칭",
         size=11, color=INK_2, line_spacing=1.85)

# Right metric card
RX = Inches(8.0); RW = Inches(4.75)
add_rect(s, RX, Inches(3.6), RW, Inches(3.3), fill=CREAM_2)
add_text(s, RX+Inches(0.35), Inches(3.8), RW-Inches(0.5), Inches(0.3),
         "WHY EVERYTIME", size=9, color=CRIMSON, bold=True, tracking=300)
add_text(s, RX+Inches(0.35), Inches(4.15), RW-Inches(0.5), Inches(1.0),
         "학생 일상 동선의 정중앙.\n광고비 0원, 도달은 최대.",
         size=16, color=INK, bold=True, line_spacing=1.25)
add_rect(s, RX+Inches(0.35), Inches(5.4), Inches(0.5), Inches(0.03), fill=CRIMSON)
add_text(s, RX+Inches(0.35), Inches(5.55), RW-Inches(0.5), Inches(1.4),
         "· 게시판은 학생들의 정보 의존도가 높음\n· '학교 근처 카공 카페' 검색 트래픽 발생\n· 추천 후기 댓글이 자연 바이럴로 작동",
         size=10, color=INK_2, line_spacing=1.85)

page_footer(s, "OUTPUT · EVERYTIME")

# ===================================================================
# SLIDE 12 — Channel Strategy Matrix
# ===================================================================
s = add_slide()
page_header(s, 10, 12, "유통 채널 전략", "CHANNEL STRATEGY")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "10", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.6),
              [("세 개의 채널, ", {"size":30,"color":INK,"bold":True,"tracking":-20}),
               ("세 가지 역할", {"size":30,"color":CRIMSON,"bold":True,"tracking":-20}),
               (".", {"size":30,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.2)

add_text(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.4),
         "각 채널의 강점이 다르므로 콘텐츠 톤과 목표를 분리해 운영합니다.",
         size=12, color=INK_2)

# Matrix - table-like rows
def matrix_row(s, y, idx, ch_kr, ch_en, role, audience, content, color_left=False):
    h = Inches(1.05)
    if color_left:
        add_rect(s, Inches(0.6), y, Inches(12.13), h, fill=CREAM_2)
    add_text(s, Inches(0.8), y+Inches(0.15), Inches(0.6), Inches(0.3),
             idx, size=10, color=CRIMSON, bold=True, tracking=200)
    add_text(s, Inches(0.8), y+Inches(0.4), Inches(2.2), Inches(0.3),
             ch_kr, size=13, color=INK, bold=True)
    add_text(s, Inches(0.8), y+Inches(0.72), Inches(2.2), Inches(0.25),
             ch_en, size=9, color=INK_3, tracking=200)
    add_text(s, Inches(3.1), y+Inches(0.3), Inches(2.6), Inches(0.6),
             role, size=11, color=INK, bold=True, line_spacing=1.4)
    add_text(s, Inches(5.85), y+Inches(0.3), Inches(3.0), Inches(0.6),
             audience, size=10, color=INK_2, line_spacing=1.5)
    add_text(s, Inches(8.95), y+Inches(0.3), Inches(3.8), Inches(0.6),
             content, size=10, color=INK_2, line_spacing=1.5)

# Header row
y_head = Inches(3.55)
add_rect(s, Inches(0.6), y_head, Inches(12.13), Inches(0.45), fill=INK)
add_text(s, Inches(0.8), y_head+Inches(0.13), Inches(2.4), Inches(0.3),
         "CHANNEL", size=10, color=CREAM, bold=True, tracking=300)
add_text(s, Inches(3.1), y_head+Inches(0.13), Inches(2.6), Inches(0.3),
         "ROLE  ·  목적", size=10, color=CREAM, bold=True, tracking=300)
add_text(s, Inches(5.85), y_head+Inches(0.13), Inches(3.0), Inches(0.3),
         "AUDIENCE  ·  타겟", size=10, color=CREAM, bold=True, tracking=300)
add_text(s, Inches(8.95), y_head+Inches(0.13), Inches(3.8), Inches(0.3),
         "CONTENT  ·  콘텐츠 톤", size=10, color=CREAM, bold=True, tracking=300)

matrix_row(s, Inches(4.05), "01", "네이버 블로그", "NAVER BLOG",
           "검색 유입 + 신뢰 구축",
           "조치원·카페 검색 의도가 있는 일반 사용자",
           "장문 1인칭 후기 · 카페 전반 소개 · 마케팅 톤", color_left=True)
matrix_row(s, Inches(5.15), "02", "에브리타임", "EVERYTIME",
           "학생 도달 + 학내 바이럴",
           "고려대·홍익대 세종캠퍼스 재학생",
           "광고 톤 X · 실사용자 후기 톤 · 학생 컨텍스트 매칭", color_left=False)
matrix_row(s, Inches(6.25), "03", "인스타그램", "INSTAGRAM",
           "지인·근거리 자연 확산",
           "본인 지인 + 조치원 근거리 거주자",
           "릴스 중심 · 비주얼 강조 · 일상 톤 큐레이션", color_left=True)

page_footer(s, "CHANNEL STRATEGY")

# ===================================================================
# SLIDE 13 — Performance Framework
# ===================================================================
s = add_slide()
page_header(s, 11, 12, "성과 측정 프레임", "PERFORMANCE FRAMEWORK")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "11", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.6),
              [("플랫폼·콘텐츠 ", {"size":28,"color":INK,"bold":True,"tracking":-20}),
               ("총합으로 측정", {"size":28,"color":CRIMSON,"bold":True,"tracking":-20}),
               ("합니다.", {"size":28,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.2)

# KPI scoreboard
add_text(s, Inches(0.6), Inches(3.3), Inches(12), Inches(0.4),
         "KPI", size=10, color=CRIMSON, bold=True, tracking=400)
add_rect(s, Inches(0.6), Inches(3.65), Inches(12.13), Inches(0.01), fill=INK)

def kpi(s, x, y, w, num, label, sub):
    add_text(s, x, y, w, Inches(1.2),
             num, size=44, color=INK, bold=True, tracking=-30, line_spacing=1.0)
    add_text(s, x, y+Inches(1.1), w, Inches(0.3),
             label, size=11, color=CRIMSON, bold=True, tracking=200)
    add_text(s, x, y+Inches(1.4), w, Inches(0.6),
             sub, size=9, color=INK_2, line_spacing=1.55)

py = Inches(3.9)
kpi(s, Inches(0.6),  py, Inches(2.6), "6", "TOTAL OUTPUTS",
    "릴스 3 + 블로그 2 + 웹페이지 1\n(추가 콘텐츠 제작 예정)")
kpi(s, Inches(3.4),  py, Inches(2.6), "4", "PLATFORMS",
    "네이버 블로그 · 에브리타임\n인스타그램 · 웹(독립 도메인)")
kpi(s, Inches(6.2),  py, Inches(2.6), "5", "METRICS TRACKED",
    "조회수 · 좋아요 · 댓글 수\n공유 수 · 체류시간")
kpi(s, Inches(9.0),  py, Inches(3.7), "TBD", "AGGREGATED SCORE",
    "현재 집계 진행 중 · 채널별 수치 통합 후 합산.\n수치 부진 시 원인 분석 + 대책 별도 정리.")

# Bottom note
add_rect(s, Inches(0.6), Inches(6.5), Inches(12.13), Inches(0.6), fill=CREAM_2)
add_text(s, Inches(0.85), Inches(6.62), Inches(11.8), Inches(0.4),
         "성과는 단일 지표가 아닌 \"콘텐츠 수 × 플랫폼 수 × 반응(좋아요·조회·공유·댓글)\"의 총합으로 평가합니다.",
         size=11, color=INK, bold=True)

page_footer(s, "PERFORMANCE FRAMEWORK")

# ===================================================================
# SLIDE 14 — Next Step
# ===================================================================
s = add_slide()
page_header(s, 12, 12, "향후 계획", "NEXT STEP")

add_text(s, Inches(0.6), Inches(1.1), Inches(2), Inches(0.4),
         "12", size=11, color=CRIMSON, bold=True, tracking=300)
add_text_runs(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.6),
              [("As-Is에서 ", {"size":30,"color":INK,"bold":True,"tracking":-20}),
               ("To-Be", {"size":30,"color":CRIMSON,"bold":True,"tracking":-20}),
               ("로, 다음 단계.", {"size":30,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.2)

# Two columns: as-is vs to-be
LX = Inches(0.6); LW = Inches(5.9)
RX = Inches(6.85); RW = Inches(5.9)

add_rect(s, LX, Inches(3.3), LW, Inches(3.6), fill=CREAM_2)
add_text(s, LX+Inches(0.4), Inches(3.5), LW-Inches(0.6), Inches(0.3),
         "AS-IS · 현재 위치", size=9, color=CRIMSON, bold=True, tracking=300)
add_text(s, LX+Inches(0.4), Inches(3.85), LW-Inches(0.6), Inches(0.5),
         "콘텐츠 자산 6종 · 채널 4종 확보", size=16, color=INK, bold=True)
add_rect(s, LX+Inches(0.4), Inches(4.45), Inches(0.5), Inches(0.03), fill=CRIMSON)
add_text(s, LX+Inches(0.4), Inches(4.6), LW-Inches(0.6), Inches(2.2),
         "· 인터랙티브 웹페이지 1종 배포 완료\n· 블로그·릴스·에브리타임 채널 운영 시작\n· AI 활용 콘텐츠 제작 워크플로 정립\n· 성과 지표 정의 및 측정 시작",
         size=11, color=INK_2, line_spacing=1.95)

add_rect(s, RX, Inches(3.3), RW, Inches(3.6), fill=INK)
add_text(s, RX+Inches(0.4), Inches(3.5), RW-Inches(0.6), Inches(0.3),
         "TO-BE · 다음 단계", size=9, color=CREAM, bold=True, tracking=300)
add_text(s, RX+Inches(0.4), Inches(3.85), RW-Inches(0.6), Inches(0.5),
         "측정 기반 채널 최적화 + 자산 확장", size=16, color=WHITE, bold=True)
add_rect(s, RX+Inches(0.4), Inches(4.45), Inches(0.5), Inches(0.03), fill=CRIMSON)
add_text(s, RX+Inches(0.4), Inches(4.6), RW-Inches(0.6), Inches(2.2),
         "· 채널별 KPI 수치 집계 → 강한 채널 더블다운\n· 부진 채널은 톤·키워드·포스팅 시간 개선\n· 콘텐츠 자산 추가 제작 (릴스·블로그 시리즈)\n· 카페 본 계정으로 콘텐츠 이관 / 위탁 운영 검토",
         size=11, color=WHITE, line_spacing=1.95)

page_footer(s, "NEXT STEP")

# ===================================================================
# SLIDE 15 — Closing
# ===================================================================
s = add_slide()

# Full bleed cream bg
add_rect(s, 0, 0, SW, SH, fill=CREAM)

# Crimson short line accent
add_rect(s, Inches(0.8), Inches(0.9), Inches(0.6), Inches(0.05), fill=CRIMSON)

add_text(s, Inches(0.8), Inches(1.15), Inches(8), Inches(0.4),
         "CLOSING", size=10, color=INK, bold=True, tracking=400)

# Big quote
add_text_runs(s, Inches(0.8), Inches(2.6), Inches(11.5), Inches(2.5),
              [("\"좋은 커피를 만드는 만큼,\n", {"size":40,"color":INK,"bold":True,"tracking":-20}),
               ("더 많은 사람이 알았으면 합니다", {"size":40,"color":CRIMSON,"bold":True,"tracking":-20}),
               (".\"", {"size":40,"color":INK,"bold":True,"tracking":-20})],
              line_spacing=1.2)

add_text(s, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.5),
         "— 카페 수작 사장님",
         size=14, color=INK_2, bold=True)

add_rect(s, Inches(0.8), Inches(6.3), Inches(12), Inches(0.01), fill=INK)

add_text(s, Inches(0.8), Inches(6.5), Inches(8), Inches(0.4),
         "이 프로젝트는 수작의 '가치'를 고객에게 전달하는 과정입니다.",
         size=12, color=INK, bold=True)
add_text(s, Inches(0.8), Inches(6.85), Inches(8), Inches(0.4),
         "Thank you.",
         size=11, color=INK_2, tracking=200, bold=True)

add_text(s, Inches(10.5), Inches(6.85), Inches(2.5), Inches(0.4),
         "CAFE SUZAK · 2026",
         size=9, color=INK_3, tracking=300, align=PP_ALIGN.RIGHT)

# ---------- Save ----------
out_path = r"C:\Users\User\Desktop\혜빈\성과보고서\콘텐츠마케팅_진행보고서.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
